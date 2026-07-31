"""
quadoa_tools.py
PowerPoint Generator for Heliostat Beam Down Analysis
"""

import os
import re
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import numpy as np
from io import BytesIO
from PIL import Image


def make_white_transparent(img_path, threshold=240):
    """Fast RGBA conversion: turn near-white pixels transparent."""
    im = Image.open(img_path).convert("RGBA")

    # w, h = im.size
    # im = im.resize((w // 2, h // 2), Image.LANCZOS)

    arr = np.asarray(im).copy()             # shape (H, W, 4)
    r, g, b, a = arr[...,0], arr[...,1], arr[...,2], arr[...,3]

    mask = (r > threshold) & (g > threshold) & (b > threshold)
    arr[mask, 3] = 0                       # zero alpha where white
    buf = BytesIO()
    Image.fromarray(arr, mode="RGBA").save(buf, format="PNG")
    buf.seek(0)
    return buf

def parse_time_from_folder(folder_name):
    """Extract time from folder name like '251008_beam_down_1_10' -> 10.0"""
    match = re.search(r'_(\d+(?:p\d+)?)$', folder_name)
    if match:
        time_str = match.group(1)
        if 'p' in time_str:
            # Handle decimal times like '12p5' -> 12.5
            parts = time_str.split('p')
            return float(parts[0]) + float(parts[1]) / 10
        else:
            return float(time_str)
    return None


def format_time(time_value):
    """Convert time value to formatted string like 10.0 -> '10:00 AM'"""
    hour = int(time_value)
    minute = int((time_value - hour) * 60)
    period = 'AM' if hour < 12 else 'PM'
    display_hour = hour if hour <= 12 else hour - 12
    display_hour = 12 if display_hour == 0 else display_hour
    return f"{display_hour}:{minute:02d} {period}"


def parse_heliostat_position(filename):
    """Extract X and Y coordinates from filename like 'HeliPosX0.0Y32.125m.png'"""
    match = re.match(r'HeliPosX(-?\d+\.?\d*)Y(-?\d+\.?\d*)m\.png', filename)
    if match:
        x = float(match.group(1))
        y = float(match.group(2))
        return (x, y)
    return None


def create_title_slide(prs, title_text):
    """Create a title slide with centered text"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title text box
    left = Inches(1)
    top = Inches(3)
    width = Inches(8)
    height = Inches(1)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = title_text
    
    # Format text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(54)
    paragraph.font.bold = True


def create_heliostat_grid_slide(prs, time_folder_path, time_str):
    """Create a 5x5 grid of heliostat images with center missing"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(8)
    height = Inches(0.5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = f"Heliostat Shapes - {time_str}"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True

    # Define grid parameters — increased by ~30%
    grid_size = 5
    aspect_ratio = 3/4  # height/width
    img_width = Inches(1.5 * 1.3)
    img_height = img_width * aspect_ratio
    spacing_x = Inches(-0.1)
    spacing_y = Inches(-0.3)
    total_width = grid_size * img_width + (grid_size - 1) * spacing_x
    total_height = grid_size * img_height + (grid_size - 1) * spacing_y
    # Margins (in inches)
    top_margin_in  = 0.8
    side_margin_in = 0.4

    start_left = Inches(side_margin_in)
    start_top  = Inches(top_margin_in)

    # Optional: auto-adjust horizontal centering if you still want it centered inside margins
    available_width = prs.slide_width - Inches(2*side_margin_in)
    offset_x = (available_width - total_width) / 2
    start_left += max(offset_x, 0)

    # Expected positions
    x_positions = [-64.25, -32.125, 0.0, 32.125, 64.25]
    y_positions = [64.25, 32.125, 0.0, -32.125, -64.25]

    # Find all heliostat images
    heliostat_files = {}
    for filename in os.listdir(time_folder_path):
        if filename.startswith('HeliPosX') and filename.endswith('.png'):
            pos = parse_heliostat_position(filename)
            if pos:
                heliostat_files[pos] = os.path.join(time_folder_path, filename)

    # Place images in grid
    for row_idx, y in enumerate(y_positions):
        for col_idx, x in enumerate(x_positions):
            if x == 0.0 and y == 0.0:
                continue

            pos = (x, y)
            if pos in heliostat_files:
                left = start_left + col_idx * (img_width + spacing_x)
                top = start_top + row_idx * (img_height + spacing_y)


                
                try:
                    # Open and preprocess image (make white transparent)
                    # img_path = heliostat_files[pos]
                    # with Image.open(img_path) as im:
                    #     im = im.convert("RGBA")
                    #     datas = im.getdata()
                    #     new_data = []
                    #     for item in datas:
                    #         # item = (r, g, b, a)
                    #         if item[0] > 240 and item[1] > 240 and item[2] > 240:
                    #             new_data.append((255, 255, 255, 0))
                    #         else:
                    #             new_data.append(item)
                    #     im.putdata(new_data)
                    #     buf = BytesIO()
                    #     im.save(buf, format="PNG")
                    #     buf.seek(0)

                    #     # Insert processed image
                    #     slide.shapes.add_picture(buf, left, top, width=img_width, height=img_height)
                    buf = make_white_transparent(heliostat_files[pos])
                    slide.shapes.add_picture(buf, left, top, width=img_width, height=img_height)
                except Exception as e:
                    print(f"Warning: Could not add image {heliostat_files[pos]}: {e}")


def create_combined_irradiance_slide(prs, irrad_folder_path, time_str):
    """Create a slide with combined irradiance images"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(8.5)
    height = Inches(0.5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = f"Combined Irradiance - {time_str}"
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True
    
    # Find the irradiance images
    encircled_energy_img = None
    irrad_combined_img = None
    
    for filename in os.listdir(irrad_folder_path):
        if filename == 'encircled_energy_combined.png':
            encircled_energy_img = os.path.join(irrad_folder_path, filename)
        elif filename.startswith('irrad_combined_contourf') and filename.endswith('.png'):
            irrad_combined_img = os.path.join(irrad_folder_path, filename)
    
    # Image dimensions - maintain aspect ratio, don't squish
    img_height = Inches(4)
    img_width_irrad = Inches(4)*2241/1936  # Square to preserve aspect ratio
    img_width_EE = Inches(4)*2102/2135  # Square to preserve aspect ratio
    vertical_center = Inches(1.5)
    
    # Add irrad_combined on the left
    if irrad_combined_img and os.path.exists(irrad_combined_img):
        left = Inches(0.25)
        try:
            slide.shapes.add_picture(
                irrad_combined_img,
                left, vertical_center,
                width=img_width_irrad,
                height=img_height
            )
        except Exception as e:
            print(f"Warning: Could not add irrad_combined image: {e}")
    
    # Add encircled_energy on the right
    if encircled_energy_img and os.path.exists(encircled_energy_img):
        left = Inches(5.5)
        try:
            slide.shapes.add_picture(
                encircled_energy_img,
                left, vertical_center,
                width=img_width_EE,
                height=img_height
            )
        except Exception as e:
            print(f"Warning: Could not add encircled_energy image: {e}")


def generate_powerpoint_of_heliostats_5by5(parent_dir, output_pptx):
    """
    Generate PowerPoint presentation from heliostat beam down analysis data.
    
    Parameters
    ----------
    parent_dir : str
        Path to the parent directory containing beam_down folders
    output_pptx : str
        Output PowerPoint filename with full path
        
    Returns
    -------
    None
    
    Examples
    --------
    >>> generate_powerpoint_of_heliostats_5by5(
    ...     r"C:\\path\\to\\winter_solstice",
    ...     r"C:\\path\\to\\winter_solstice\\combined_results.pptx"
    ... )
    """
    parent_path = Path(parent_dir)
    
    if not parent_path.exists():
        raise ValueError(f"Directory does not exist: {parent_dir}")
    
    # Get parent folder name for title
    title_text = parent_path.name.replace('_', ' ').title()
    
    # Find all beam_down folders and sort by time
    time_folders = []
    for item in parent_path.iterdir():
        if item.is_dir() and 'beam_down' in item.name.lower():
            time_value = parse_time_from_folder(item.name)
            if time_value is not None:
                time_folders.append((time_value, item))
    
    # Sort by time
    time_folders.sort(key=lambda x: x[0])
    
    if not time_folders:
        raise ValueError(f"No beam_down folders found in {parent_dir}")
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Create title slide
    create_title_slide(prs, title_text)
    
    # Process each time folder
    for time_value, time_folder in time_folders:
        time_str = format_time(time_value)
        print(f"Processing {time_folder.name} ({time_str})...")
        
        # Create heliostat grid slide
        create_heliostat_grid_slide(prs, time_folder, time_str)
        
        # Create combined irradiance slide
        irrad_folder = time_folder / 'irrad'
        if irrad_folder.exists():
            create_combined_irradiance_slide(prs, irrad_folder, time_str)
        else:
            print(f"Warning: irrad folder not found in {time_folder.name}")
    
    # Save presentation
    prs.save(output_pptx)
    print(f"\nPresentation saved to: {output_pptx}")